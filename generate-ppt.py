#!/usr/bin/env python3
"""산업동향 보고서 PPT 생성 - 템플릿 없이 자체 디자인
깔끔한 비즈니스 스타일, 상부 보고용
"""

import json, sys, os, re
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt, Cm, Inches, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

FONT = "맑은 고딕"

# 컬러 팔레트
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY = RGBColor(0x0F, 0x1A, 0x30)
BLUE = RGBColor(0x2E, 0x75, 0xB6)
LIGHT_BLUE = RGBColor(0xD6, 0xE4, 0xF0)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x76, 0x76, 0x76)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
BORDER_GRAY = RGBColor(0xD9, 0xD9, 0xD9)

W = Cm(25.4)
H = Cm(19.05)

ACCENT_COLORS = [
    RGBColor(0x2E, 0x75, 0xB6),
    RGBColor(0xE8, 0x6C, 0x00),
    RGBColor(0x54, 0x8B, 0x54),
    RGBColor(0x8B, 0x5C, 0xF6),
    RGBColor(0xC0, 0x39, 0x2B),
]


def trunc(text, n=60):
    text = text.strip()
    if len(text) <= n:
        return text
    for end in ["됩니다", "입니다", "습니다", "겠습니다", "니다"]:
        idx = text.find(end, n // 3)
        if 0 < idx < n + 20:
            return text[:idx + len(end)]
    for sep in [". ", ", ", "; "]:
        idx = text.rfind(sep, n // 3, n + 15)
        if idx > 0:
            return text[:idx + 1]
    return text[:n] + "…"


def rect(slide, left, top, width, height, color, alpha=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def text_box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    return tb.text_frame


def add_para(tf, text, size=11, bold=False, color=BLACK, spacing=0, align=PP_ALIGN.LEFT):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(spacing)
    p.alignment = align
    p.line_spacing = Pt(size * 1.5)
    return p


def add_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text
    for p in tf.paragraphs:
        p.font.name = FONT
        p.font.size = Pt(12)


def footer_bar(slide):
    """하단 바: 좌=로고텍스트, 우=부서"""
    rect(slide, Cm(0), Cm(18.2), W, Cm(0.85), NAVY)
    tf = text_box(slide, Cm(0.8), Cm(18.25), Cm(10), Cm(0.7))
    add_para(tf, "ISU GROUP", size=8, bold=True, color=WHITE)
    tf2 = text_box(slide, Cm(15), Cm(18.25), Cm(10), Cm(0.7))
    add_para(tf2, "(주)이수 기획팀", size=8, color=WHITE, align=PP_ALIGN.RIGHT)


def top_bar(slide, title, section_num=None):
    """상단 제목 바"""
    rect(slide, Cm(0), Cm(0), W, Cm(1.4), NAVY)
    label = f"{section_num}. {title}" if section_num else title
    tf = text_box(slide, Cm(0.5), Cm(0.15), Cm(22), Cm(1.1))
    add_para(tf, label, size=16, bold=True, color=WHITE)
    # 우측 악센트 바
    rect(slide, Cm(24.4), Cm(0), Cm(1.0), Cm(1.4), ORANGE)


def head_message(slide, text):
    """헤드메시지: 2줄 100자"""
    rect(slide, Cm(0.8), Cm(1.8), Cm(23.8), Cm(1.6), LIGHT_BLUE)
    tf = text_box(slide, Cm(1.2), Cm(1.9), Cm(23.0), Cm(1.4))
    add_para(tf, trunc(text, 100), size=13, bold=True, color=NAVY)


def strip_num(title):
    """제목에서 선행 번호 제거: '2. 산업 구조' → '산업 구조'"""
    m = re.match(r'^[\d]+[\.\)]\s*', title)
    return title[m.end():] if m else title


def parse_sections(text):
    secs = []
    cur = None
    sub = None
    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("## "):
            if cur: secs.append(cur)
            cur = {"title": strip_num(line[3:].strip()), "subs": [], "content": []}
            sub = None
        elif line.startswith("### "):
            if cur:
                sub = {"title": line[4:].strip(), "content": []}
                cur["subs"].append(sub)
        elif line.startswith("# "):
            continue
        elif cur:
            clean = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
            if line.startswith("- ") or line.startswith("* "):
                item = {"text": clean[2:], "level": 1}
            elif line.startswith("  - ") or line.startswith("  * "):
                item = {"text": clean[4:], "level": 2}
            else:
                item = {"text": clean, "level": 0}
            if sub:
                sub["content"].append(item)
            else:
                cur["content"].append(item)
    if cur:
        secs.append(cur)
    return secs


def get_head(content, subs):
    for it in content:
        if it["level"] == 0 and len(it["text"]) > 20:
            return trunc(it["text"], 100)
    for s in subs:
        for it in s.get("content", []):
            if it["level"] == 0 and len(it["text"]) > 20:
                return trunc(it["text"], 100)
    return None


def fallback_script(title, items):
    lines = [f"[{title}]", ""]
    for it in items:
        if it.get("level", 0) == 0 and len(it["text"]) > 15:
            lines.append(it["text"])
            lines.append("")
    sub = [it for it in items if it.get("level") == 1]
    if sub:
        for s in sub[:5]:
            lines.append(f"  - {s['text']}")
    return "\n".join(lines)


# ── 슬라이드 ──

def slide_cover(prs, label, date, sub_field):
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 배경
    rect(sl, Cm(0), Cm(0), W, H, DARK_NAVY)

    # 좌측 악센트 바
    rect(sl, Cm(0), Cm(0), Cm(0.5), H, BLUE)

    # 오렌지 악센트 라인
    rect(sl, Cm(3.0), Cm(7.5), Cm(3.0), Cm(0.15), ORANGE)

    # 제목
    tf = text_box(sl, Cm(3.0), Cm(8.0), Cm(19), Cm(3.0))
    add_para(tf, f"{label}", size=36, bold=True, color=WHITE)
    add_para(tf, "산업동향 초안", size=28, color=RGBColor(0xA0, 0xB4, 0xCC), spacing=8)

    # 소제목
    if sub_field:
        tf2 = text_box(sl, Cm(3.0), Cm(11.5), Cm(19), Cm(1.5))
        add_para(tf2, sub_field, size=16, color=RGBColor(0x80, 0x99, 0xB3))

    # 하단 정보
    tf3 = text_box(sl, Cm(3.0), Cm(15.5), Cm(19), Cm(2.0))
    add_para(tf3, date, size=14, color=GRAY)
    add_para(tf3, "(주)이수 기획팀", size=14, bold=True, color=RGBColor(0xA0, 0xB4, 0xCC), spacing=4)

    sf = f" 중 {sub_field} 분야" if sub_field else ""
    add_notes(sl, f"안녕하십니까. 지금부터 {label}{sf} 산업동향에 대해 보고드리겠습니다.\n\n본 보고서는 국내외 뉴스, 학술논문, 정책자료 등을 종합 분석하여 작성되었습니다.")


def slide_exec_summary(prs, sections):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    top_bar(sl, "Executive Summary")
    footer_bar(sl)

    summaries = []
    for sec in sections:
        if "executive" in sec["title"].lower() or "summary" in sec["title"].lower():
            continue
        hd = get_head(sec["content"], sec.get("subs", []))
        bullets = []
        for it in sec["content"] + [it for s in sec.get("subs",[]) for it in s.get("content",[])]:
            if it["level"] == 1:
                bullets.append(trunc(it["text"], 50))
        summaries.append({"title": sec["title"], "head": hd or sec["title"], "bullets": bullets[:2]})
        if len(summaries) >= 3:
            break

    y = Cm(2.0)
    for i, sm in enumerate(summaries):
        color = ACCENT_COLORS[i % len(ACCENT_COLORS)]
        box_h = Cm(4.6)

        # 좌측 컬러 바
        rect(sl, Cm(1.0), y + Cm(0.3), Cm(0.25), box_h - Cm(0.6), color)

        # 넘버
        sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1.6), y + Cm(0.5), Cm(1.0), Cm(1.0))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        ntf = sh.text_frame
        ntf.margin_left = ntf.margin_right = ntf.margin_top = ntf.margin_bottom = Cm(0)
        np = ntf.paragraphs[0]
        np.text = str(i + 1)
        np.font.name = FONT; np.font.size = Pt(14); np.font.bold = True
        np.font.color.rgb = WHITE; np.alignment = PP_ALIGN.CENTER

        # 제목
        tf = text_box(sl, Cm(3.0), y + Cm(0.3), Cm(21), Cm(0.8))
        add_para(tf, sm["title"], size=14, bold=True, color=NAVY)

        # 요약
        tf2 = text_box(sl, Cm(3.0), y + Cm(1.2), Cm(21), Cm(0.8))
        add_para(tf2, trunc(sm["head"], 80), size=11, color=BLACK)

        # 불릿
        tf3 = text_box(sl, Cm(3.5), y + Cm(2.2), Cm(20.5), Cm(2.0))
        for b in sm["bullets"]:
            add_para(tf3, f"▸  {b}", size=10, color=GRAY, spacing=3)

        y += box_h + Cm(0.5)

    script = "먼저 핵심 내용을 요약드리겠습니다.\n\n"
    for sm in summaries:
        script += f"{sm['title']}과 관련하여, {sm['head']}\n\n"
    add_notes(sl, script)


def slide_toc(prs, sections):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    top_bar(sl, "목 차")
    footer_bar(sl)

    for i, sec in enumerate(sections):
        y = Cm(2.8) + Cm(i * 2.2)
        if y > Cm(16.5):
            break
        color = ACCENT_COLORS[i % len(ACCENT_COLORS)]

        # 넘버 원
        sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, Cm(3.5), y, Cm(1.0), Cm(1.0))
        sh.fill.solid()
        sh.fill.fore_color.rgb = color
        sh.line.fill.background()
        ntf = sh.text_frame
        ntf.margin_left = ntf.margin_right = ntf.margin_top = ntf.margin_bottom = Cm(0)
        np = ntf.paragraphs[0]
        np.text = str(i + 1)
        np.font.name = FONT; np.font.size = Pt(14); np.font.bold = True
        np.font.color.rgb = WHITE; np.alignment = PP_ALIGN.CENTER

        # 제목
        tf = text_box(sl, Cm(5.2), y + Cm(0.1), Cm(18), Cm(0.9))
        add_para(tf, sec["title"], size=16, bold=True, color=NAVY)

        # 구분선
        if i < len(sections) - 1:
            line = sl.shapes.add_connector(1, Cm(5.2), y + Cm(1.7), Cm(21), y + Cm(1.7))
            line.line.color.rgb = BORDER_GRAY
            line.line.width = Pt(0.5)

    sc = "보고서 구성은 다음과 같습니다.\n\n"
    for i, s in enumerate(sections, 1):
        sc += f"{i}번째로 {s['title']},\n"
    sc += "\n순서로 말씀드리겠습니다."
    add_notes(sl, sc)


def slide_two_col(prs, title, head, lt, li, rt, ri, sec_num=None, script=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    top_bar(sl, title, sec_num)
    footer_bar(sl)

    y_start = Cm(1.8)
    if head:
        head_message(sl, head)
        y_start = Cm(3.8)

    # 좌측 소제목
    rect(sl, Cm(1.0), y_start, Cm(11.2), Cm(1.0), LIGHT_GRAY)
    tf_lt = text_box(sl, Cm(1.3), y_start + Cm(0.1), Cm(10.6), Cm(0.8))
    add_para(tf_lt, lt, size=13, bold=True, color=NAVY)

    # 좌측 구분선
    line1 = sl.shapes.add_connector(1, Cm(1.0), y_start + Cm(1.0), Cm(12.2), y_start + Cm(1.0))
    line1.line.color.rgb = BLUE
    line1.line.width = Pt(1.5)

    # 좌측 본문
    ly = y_start + Cm(1.3)
    for i, it in enumerate(li[:5]):
        lv = it.get("level", 0)
        tx = trunc(it.get("text", ""), 50)
        color = ACCENT_COLORS[i % len(ACCENT_COLORS)]

        if lv <= 1:
            # 아이콘
            dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1.3), ly + Cm(0.15), Cm(0.35), Cm(0.35))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

            tf = text_box(sl, Cm(2.0), ly, Cm(10.0), Cm(1.2))
            add_para(tf, tx, size=11, bold=(lv == 0), color=BLACK)
            ly += Cm(1.5)
        else:
            tf = text_box(sl, Cm(2.3), ly, Cm(9.7), Cm(0.8))
            add_para(tf, f"→ {tx}", size=10, color=GRAY)
            ly += Cm(1.1)

    # 우측 소제목
    rect(sl, Cm(13.2), y_start, Cm(11.2), Cm(1.0), LIGHT_GRAY)
    tf_rt = text_box(sl, Cm(13.5), y_start + Cm(0.1), Cm(10.6), Cm(0.8))
    add_para(tf_rt, rt, size=13, bold=True, color=NAVY)

    # 우측 구분선
    line2 = sl.shapes.add_connector(1, Cm(13.2), y_start + Cm(1.0), Cm(24.4), y_start + Cm(1.0))
    line2.line.color.rgb = BLUE
    line2.line.width = Pt(1.5)

    # 우측 본문
    ry = y_start + Cm(1.3)
    for i, it in enumerate(ri[:5]):
        lv = it.get("level", 0)
        tx = trunc(it.get("text", ""), 50)
        color = ACCENT_COLORS[i % len(ACCENT_COLORS)]

        if lv <= 1:
            dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Cm(13.5), ry + Cm(0.15), Cm(0.35), Cm(0.35))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

            tf = text_box(sl, Cm(14.2), ry, Cm(10.0), Cm(1.2))
            add_para(tf, tx, size=11, bold=(lv == 0), color=BLACK)
            ry += Cm(1.5)
        else:
            tf = text_box(sl, Cm(14.5), ry, Cm(9.7), Cm(0.8))
            add_para(tf, f"→ {tx}", size=10, color=GRAY)
            ry += Cm(1.1)

    if script:
        add_notes(sl, script)
    else:
        add_notes(sl, fallback_script(title, li + ri))


def slide_single(prs, title, head, items, sec_num=None, script=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    top_bar(sl, title, sec_num)
    footer_bar(sl)

    y = Cm(1.8)
    if head:
        head_message(sl, head)
        y = Cm(3.8)

    for i, it in enumerate(items[:6]):
        lv = it.get("level", 0)
        tx = trunc(it.get("text", ""), 75)
        color = ACCENT_COLORS[i % len(ACCENT_COLORS)]

        if lv == 0:
            dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, Cm(1.5), y + Cm(0.15), Cm(0.4), Cm(0.4))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

            tf = text_box(sl, Cm(2.3), y, Cm(21.5), Cm(1.2))
            add_para(tf, tx, size=13, bold=True, color=NAVY)
            y += Cm(1.8)
        elif lv == 1:
            tf = text_box(sl, Cm(2.8), y, Cm(21.0), Cm(1.0))
            add_para(tf, f"▸  {tx}", size=11, color=BLACK)
            y += Cm(1.4)
        else:
            tf = text_box(sl, Cm(3.3), y, Cm(20.5), Cm(0.8))
            add_para(tf, f"→ {tx}", size=10, color=GRAY)
            y += Cm(1.1)

    if script:
        add_notes(sl, script)
    else:
        add_notes(sl, fallback_script(title, items))


def slide_end(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, Cm(0), Cm(0), W, H, DARK_NAVY)
    rect(sl, Cm(0), Cm(0), Cm(0.5), H, BLUE)

    rect(sl, Cm(8.0), Cm(8.5), Cm(3.0), Cm(0.15), ORANGE)

    tf = text_box(sl, Cm(3.0), Cm(9.0), Cm(19.4), Cm(3.0))
    add_para(tf, "감사합니다", size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tf2 = text_box(sl, Cm(3.0), Cm(14.0), Cm(19.4), Cm(1.0))
    add_para(tf2, "(주)이수 기획팀", size=14, color=RGBColor(0xA0, 0xB4, 0xCC), align=PP_ALIGN.CENTER)

    add_notes(sl, "이상으로 산업동향 보고를 마치겠습니다.\n\n궁금하신 사항이나 추가 논의가 필요한 부분이 있으시면 말씀해 주십시오.\n\n감사합니다.")


def generate_ppt(data):
    prs = Presentation()
    prs.slide_width = Emu(int(W))
    prs.slide_height = Emu(int(H))

    label = data.get("label", "산업")
    sub_field = data.get("subField", None)
    date = data.get("generatedAt", "")[:10]
    report = data.get("report", "")
    refs = data.get("references", [])
    scripts = data.get("scripts", {})

    sections = parse_sections(report)

    exec_sec = None
    body_secs = []
    for sec in sections:
        if "executive" in sec["title"].lower() or "summary" in sec["title"].lower():
            exec_sec = sec
        else:
            body_secs.append(sec)

    # 1. 표지
    slide_cover(prs, label, date, sub_field)

    # 2. Executive Summary
    slide_exec_summary(prs, body_secs)

    # 3. 목차
    if body_secs:
        slide_toc(prs, body_secs)

    # 4. 본문
    for si, sec in enumerate(body_secs):
        sc = None
        for sk, sv in scripts.items():
            if sk in sec["title"] or sec["title"] in sk:
                sc = sv
                break

        subs = sec.get("subs", [])
        content = sec.get("content", [])
        head = get_head(content, subs)
        # head를 제거하되, 남는 항목이 있을 때만
        display_content = content
        if head and len(content) > 1:
            display_content = [it for it in content if trunc(it.get("text", ""), 100) != head]

        sec_num = si + 1

        if subs and len(subs) >= 2:
            i = 0
            sn = 0
            while i < len(subs):
                ls = subs[i]
                rs = subs[i + 1] if i + 1 < len(subs) else None
                lt, li = ls["title"], ls.get("content", [])
                if rs:
                    rt, ri = rs["title"], rs.get("content", [])
                    i += 2
                else:
                    # 홀수 sub: 우측에 전체 섹션 content 활용
                    rt = "핵심 시사점"
                    ri = display_content[:4] if display_content else li[:3]
                    i += 1
                t = sec["title"] if sn == 0 else f"{sec['title']} (계속)"
                hm = head if sn == 0 else get_head(li + ri, [])
                s = sc if sn == 0 else None
                n = sec_num if sn == 0 else None
                slide_two_col(prs, t, hm, lt, li, rt, ri, sec_num=n, script=s)
                sn += 1
        elif len(display_content) > 6:
            mid = len(display_content) // 2
            slide_two_col(prs, sec["title"], head, "주요 동향", display_content[:mid], "핵심 시사점", display_content[mid:], sec_num=sec_num, script=sc)
        else:
            # 단일 컬럼: 모든 콘텐츠 표시 (sub 포함)
            all_items = list(display_content)
            for s in subs:
                all_items.append({"text": s["title"], "level": 0})
                all_items.extend(s.get("content", []))
            # content가 비어있으면 원본 텍스트를 문장 단위로 분해
            if not all_items:
                for it in content:
                    sentences = [s.strip() for s in it["text"].split(". ") if s.strip()]
                    for sent in sentences:
                        all_items.append({"text": sent + ("." if not sent.endswith(".") else ""), "level": 1})
            if not all_items:
                all_items = [{"text": sec["title"], "level": 0}]
            slide_single(prs, sec["title"], head, all_items, sec_num=sec_num, script=sc)

    # 5. 참고자료
    if refs:
        dom = [r for r in refs if r.get("type") != "global"][:5]
        glb = [r for r in refs if r.get("type") == "global"][:5]
        li = [{"text": trunc(r.get("title", ""), 45), "level": 1} for r in dom]
        ri = [{"text": trunc(r.get("title", ""), 45), "level": 1} for r in (glb or refs[5:10])]
        slide_two_col(prs, "참고자료", None, "국내 자료", li, "글로벌 자료", ri)

    # 6. 끝
    slide_end(prs)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


if __name__ == "__main__":
    data = json.load(sys.stdin)
    sys.stdout.buffer.write(generate_ppt(data))
